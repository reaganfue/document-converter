"""
converters/ppt_converter.py — PPTX 來源轉換器

支援：PPTX → PDF / HTML / TXT / IMG
首選 PPTX→PDF：python-pptx 抽元素 + Pillow 繪製 → 合成 PDF
降級 PPTX→PDF：逐頁 PNG → Pillow 合成 PDF（同首選邏輯）
備注：不支援動畫、影片、複雜 OLE 物件
"""

import html as html_module
import logging
import platform
import zipfile
from pathlib import Path

from .base import BaseConverter
from .exceptions import ConversionError, UnsupportedFormatError
from .utils import build_minimal_html, ensure_output_dir

logger = logging.getLogger(__name__)

# 投影片渲染解析度（點/英吋）
_SLIDE_DPI = 96
# 標準投影片比例 10×7.5 英吋（寬×高）
_SLIDE_WIDTH_PX = 960
_SLIDE_HEIGHT_PX = 720


class PPTConverter(BaseConverter):
    """PPTX 格式來源轉換器。

    支援目標格式：pdf / html / txt / image
    """

    source_format = "pptx"
    supported_targets = ["pdf", "html", "txt", "image"]

    def convert(self, input_path: Path, output_path: Path, target_format: str) -> None:
        """執行 PPTX 轉換。

        Args:
            input_path: 來源 PPTX 路徑
            output_path: 目標檔案路徑
            target_format: "pdf" | "html" | "txt" | "image"

        Raises:
            UnsupportedFormatError: target_format 不在支援清單
            ConversionError: 轉換失敗
        """
        if not self.supports(target_format):
            raise UnsupportedFormatError(
                f"PPTConverter 不支援目標格式: {target_format}"
            )

        ensure_output_dir(output_path)
        self._log_primary_attempt(target_format)

        dispatch_map = {
            "pdf": self._to_pdf,
            "html": self._to_html,
            "txt": self._to_txt,
            "image": self._to_image,
        }
        dispatch_map[target_format](input_path, output_path)

    # ─── 各目標格式實作 ────────────────────────────────────────────

    def _to_pdf(self, input_path: Path, output_path: Path) -> None:
        """PPTX → PDF：Pillow 繪製每頁 → 合成 PDF。

        首選：Pillow 圖像模式（Slides 繪製 + 合成 PDF）
        降級：抽文字 → HTML → weasyprint
        """
        import shutil
        primary_error: Exception | None = None
        png_paths: list[Path] = []
        try:
            png_paths = self._render_slides_to_pngs(input_path)
            self._pngs_to_pdf(png_paths, output_path)
            logger.info("PPTX→PDF 首選成功: Pillow 渲染")
            return
        except Exception as exc:
            primary_error = exc
            self._log_fallback_attempt("pdf", exc)
        finally:
            # 清理 _render_slides_to_pngs 建立的 persistent_dir
            if png_paths:
                parent = png_paths[0].parent
                shutil.rmtree(str(parent), ignore_errors=True)

        # 降級：文字抽取 → HTML → 多層渲染降級鏈
        try:
            from .utils import temp_file
            from .pdf_renderer import get_renderer

            html_content = self._to_html_string(input_path)
            with temp_file(suffix=".html") as html_path:
                html_path.write_text(html_content, encoding="utf-8")
                get_renderer().render_html_to_pdf(html_path, output_path)
            logger.info("PPTX→PDF 降級成功: pdf_renderer")
        except Exception as fallback_error:
            self._raise_conversion_error("pdf", primary_error, fallback_error)

    def _to_html(self, input_path: Path, output_path: Path) -> None:
        """PPTX → HTML：python-pptx 抽 shapes → 組 HTML。"""
        try:
            html_content = self._to_html_string(input_path)
            output_path.write_text(html_content, encoding="utf-8")
            logger.info("PPTX→HTML 成功")
        except Exception as exc:
            raise ConversionError("PPTX→HTML 轉換失敗", source_error=exc) from exc

    def _to_txt(self, input_path: Path, output_path: Path) -> None:
        """PPTX → TXT：python-pptx 每 shape .text 串接。"""
        try:
            from pptx import Presentation

            prs = Presentation(str(input_path))
            lines: list[str] = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                lines.append(f"\n=== 第 {slide_num} 頁 ===")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)

            output_path.write_text("\n".join(lines), encoding="utf-8")
            logger.info("PPTX→TXT 成功")
        except Exception as exc:
            raise ConversionError("PPTX→TXT 轉換失敗", source_error=exc) from exc

    def _to_image(self, input_path: Path, output_path: Path) -> None:
        """PPTX → IMG：Pillow 繪製，單頁輸出 PNG，多頁輸出 ZIP。"""
        import shutil
        png_paths: list[Path] = []
        try:
            png_paths = self._render_slides_to_pngs(input_path)

            if len(png_paths) == 1:
                shutil.copy(str(png_paths[0]), str(output_path))
            else:
                zip_output = output_path.with_suffix(".zip")
                with zipfile.ZipFile(str(zip_output), "w", zipfile.ZIP_DEFLATED) as zf:
                    for png_path in png_paths:
                        zf.write(str(png_path), png_path.name)
                if output_path != zip_output:
                    if output_path.exists():
                        output_path.unlink()
                    zip_output.rename(output_path)

            logger.info("PPTX→IMG 成功（%d 頁）", len(png_paths))
        except Exception as exc:
            raise ConversionError("PPTX→IMG 轉換失敗", source_error=exc) from exc
        finally:
            # 清理 _render_slides_to_pngs 建立的 persistent_dir（整個父目錄）
            if png_paths:
                parent = png_paths[0].parent
                try:
                    shutil.rmtree(str(parent), ignore_errors=True)
                except OSError:
                    pass

    # ─── 共用輔助 ──────────────────────────────────────────────────

    def _resolve_font(self, font_size_px: int):
        """依作業系統選擇最佳可用字型，回退至 Pillow 預設。

        Args:
            font_size_px: 字型像素大小

        Returns:
            ImageFont instance（TrueType 或 default）
        """
        from PIL import ImageFont

        if platform.system() == "Windows":
            candidates = [
                "C:/Windows/Fonts/msjh.ttc",   # 微軟正黑（CJK 首選）
                "C:/Windows/Fonts/msyh.ttc",   # 微軟雅黑
                "C:/Windows/Fonts/arial.ttf",  # 英文備援
            ]
        else:
            candidates = [
                "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",  # Linux CJK
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/System/Library/Fonts/Helvetica.ttc",  # macOS
            ]

        for path in candidates:
            try:
                return ImageFont.truetype(path, font_size_px)
            except (IOError, OSError):
                continue

        return ImageFont.load_default()

    def _render_slides_to_pngs(self, input_path: Path) -> list[Path]:
        """用 Pillow 繪製每張投影片為 PNG，回傳路徑清單。

        使用 tempfile.TemporaryDirectory context manager 確保例外時不洩漏臨時目錄
        （修復 W-PERF-004）。
        """
        import tempfile
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image, ImageDraw

        prs = Presentation(str(input_path))

        # 取投影片實際大小（EMU → pixels）
        slide_width_emu = prs.slide_width or Emu(9144000)   # 預設 10 英吋
        slide_height_emu = prs.slide_height or Emu(6858000)  # 預設 7.5 英吋
        emu_per_px = 914400 / _SLIDE_DPI  # 1 英吋 = 914400 EMU

        width_px = int(slide_width_emu / emu_per_px)
        height_px = int(slide_height_emu / emu_per_px)

        # 0 頁的 PPT 提前拋錯，避免後續建立 persistent_dir 後無法清理（tmpdir 洩漏修復）
        if len(prs.slides) == 0:
            raise ConversionError("PPT 沒有任何投影片，無法轉換")

        png_paths: list[Path] = []

        # TemporaryDirectory 在 exception 時自動清理（W-PERF-004）
        with tempfile.TemporaryDirectory(prefix="pptx_render_") as tmp_dir_str:
            tmp_dir = Path(tmp_dir_str)

            for slide_num, slide in enumerate(prs.slides, start=1):
                img = Image.new("RGB", (width_px, height_px), color=(255, 255, 255))
                draw = ImageDraw.Draw(img)

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue

                    # 取形狀位置與大小（EMU → px）
                    left_px = int(shape.left / emu_per_px) if shape.left else 10
                    top_px = int(shape.top / emu_per_px) if shape.top else 10

                    y_offset = top_px
                    for para in shape.text_frame.paragraphs:
                        text = para.text
                        if not text.strip():
                            y_offset += 12
                            continue

                        # 判斷字型大小（pptx pt → px）
                        font_size_pt = 18
                        if para.runs:
                            for run in para.runs:
                                if run.font.size:
                                    from pptx.util import Pt
                                    font_size_pt = int(run.font.size.pt)
                                    break

                        font_size_px = max(10, int(font_size_pt * _SLIDE_DPI / 72))
                        font = self._resolve_font(font_size_px)

                        # 顏色
                        text_color = (0, 0, 0)
                        if para.runs and para.runs[0].font.color and para.runs[0].font.color.type:
                            try:
                                rgb = para.runs[0].font.color.rgb
                                text_color = (rgb.r, rgb.g, rgb.b)
                            except Exception:
                                pass

                        draw.text((left_px, y_offset), text, fill=text_color, font=font)
                        bbox = font.getbbox(text) if hasattr(font, "getbbox") else (0, 0, 0, font_size_px + 4)
                        line_height = bbox[3] - bbox[1] + 4
                        y_offset += line_height

                png_path = tmp_dir / f"slide_{slide_num:04d}.png"
                img.save(str(png_path), "PNG")
                img.close()
                # 收集路徑供呼叫者使用（PNG 已存盤，tmp_dir 仍存在）
                png_paths.append(png_path)

            # 在 with 區塊結束前，呼叫者必須已完成對 png_paths 的使用。
            # _pngs_to_pdf 和 _to_image 在此 context 外呼叫，需在 context 內回傳資料。
            # 因此將 PNG 複製到呼叫者可存取的路徑，再讓 TemporaryDirectory 自行清理。
            import shutil as _shutil
            import tempfile as _tempfile
            persistent_dir = Path(_tempfile.mkdtemp(prefix="pptx_out_"))
            persistent_paths: list[Path] = []
            for src in png_paths:
                dst = persistent_dir / src.name
                _shutil.copy2(str(src), str(dst))
                persistent_paths.append(dst)

        return persistent_paths

    def _pngs_to_pdf(self, png_paths: list[Path], output_path: Path) -> None:
        """用 Pillow 將多張 PNG 合成 PDF。

        使用串流（生成器）方式逐頁處理，避免全量載入至記憶體（W-PERF-002）。
        Pillow save_all + append_images 接受生成器，每頁處理後即可釋放。
        """
        from PIL import Image

        if not png_paths:
            raise ConversionError("沒有任何投影片可轉換")

        if len(png_paths) == 1:
            with Image.open(str(png_paths[0])) as img:
                img.convert("RGB").save(str(output_path), "PDF", resolution=_SLIDE_DPI)
            return

        # 串流方式：第一頁開啟後寫入，其餘用生成器逐頁 append
        def _remaining_images():
            for path in png_paths[1:]:
                with Image.open(str(path)) as img:
                    yield img.convert("RGB")

        with Image.open(str(png_paths[0])) as first_img:
            first_rgb = first_img.convert("RGB")
            first_rgb.save(
                str(output_path),
                "PDF",
                resolution=_SLIDE_DPI,
                save_all=True,
                append_images=_remaining_images(),
            )

    def _to_html_string(self, input_path: Path) -> str:
        """python-pptx 抽 shapes → HTML（每 slide 一個 <section>）。"""
        from pptx import Presentation

        prs = Presentation(str(input_path))
        sections: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = [f"<h2>第 {slide_num} 頁</h2>"]

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        escaped = html_module.escape(text)
                        slide_parts.append(f"<p>{escaped}</p>")

            sections.append(
                f'<section class="slide">\n'
                + "\n".join(slide_parts)
                + "\n</section>"
            )

        body = "\n<hr>\n".join(sections)
        extra_style = ".slide { margin: 2em 0; padding: 1em; border: 1px solid #ddd; }"
        return build_minimal_html(body, title="投影片轉換結果", extra_style=extra_style)
