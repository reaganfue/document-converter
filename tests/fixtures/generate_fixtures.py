"""
tests/fixtures/generate_fixtures.py — 程式化生成測試 fixtures

不使用任何版權素材，所有 fixtures 純程式生成。
呼叫 generate_all() 建立所有測試檔案。
"""

from pathlib import Path

FIXTURES_DIR = Path(__file__).parent


def generate_sample_pdf() -> Path:
    """建立測試用 PDF。

    優先使用 PyMuPDF（fitz）生成高品質 PDF；
    若未安裝則回退到最小化合法 PDF 結構。
    """
    output_path = FIXTURES_DIR / "sample.pdf"

    try:
        import fitz

        doc = fitz.open()

        # 第 1 頁：標題 + 段落
        page1 = doc.new_page(width=595, height=842)  # A4
        page1.insert_text(
            (72, 100),
            "Test Document (Testing Only)",
            fontsize=24,
            color=(0, 0, 0),
        )
        page1.insert_text(
            (72, 150),
            "Page 1 - Test Content",
            fontsize=14,
        )
        paragraph = (
            "This is a sample paragraph in English for testing conversion.\n"
            "The document contains basic text for testing purposes."
        )
        for i, line in enumerate(paragraph.split("\n")):
            page1.insert_text((72, 200 + i * 30), line, fontsize=11)

        # 第 2 頁：簡單表格（文字模擬）
        page2 = doc.new_page(width=595, height=842)
        page2.insert_text((72, 80), "Page 2 - Table Example", fontsize=14)

        headers = ["Name", "Format", "Size"]
        rows = [
            ["Document A", "PDF", "1.2 MB"],
            ["Document B", "DOCX", "0.8 MB"],
            ["Document C", "HTML", "0.3 MB"],
        ]

        col_x = [72, 200, 330]
        page2.insert_text((72, 120), " | ".join(headers), fontsize=12)
        page2.draw_line((72, 135), (450, 135), color=(0, 0, 0))

        for row_idx, row in enumerate(rows):
            for col_idx, cell in enumerate(row):
                page2.insert_text(
                    (col_x[col_idx], 150 + row_idx * 25),
                    cell,
                    fontsize=11,
                )

        doc.save(str(output_path))
        doc.close()
    except ImportError:
        # 回退：建立最小化合法 PDF（純 ASCII 內容，無需 CJK 字型）
        _write_minimal_pdf(output_path)

    return output_path


def _write_minimal_pdf(output_path: Path) -> None:
    """建立最小化合法 PDF（不依賴任何第三方庫）。"""
    pdf_body = (
        b"%PDF-1.4\n"
        b"1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n"
        b"2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
        b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842]\n"
        b"   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\nendobj\n"
        b"4 0 obj\n<< /Length 185 >>\nstream\n"
        b"BT\n"
        b"/F1 18 Tf\n72 750 Td\n"
        b"(Test Document - Converter Test Fixture) Tj\n"
        b"0 -30 Td\n(This is page 1 of the test PDF.) Tj\n"
        b"0 -30 Td\n(It contains plain ASCII text for testing.) Tj\n"
        b"0 -30 Td\n(Name | Format | Size) Tj\n"
        b"0 -25 Td\n(Document A | PDF | 1.2 MB) Tj\n"
        b"0 -25 Td\n(Document B | DOCX | 0.8 MB) Tj\n"
        b"ET\n"
        b"endstream\nendobj\n"
        b"5 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000295 00000 n \n"
        b"0000000532 00000 n \n"
        b"trailer\n<< /Size 6 /Root 1 0 R >>\n"
        b"startxref\n613\n%%EOF\n"
    )
    output_path.write_bytes(pdf_body)


def generate_sample_docx() -> Path:
    """用 python-docx 建立含標題、正文、項目符號的 DOCX。"""
    from docx import Document
    from docx.shared import Pt

    output_path = FIXTURES_DIR / "sample.docx"
    doc = Document()

    doc.add_heading("測試 Word 文件", level=1)
    doc.add_heading("第一節 — 簡介", level=2)
    doc.add_paragraph(
        "這是一份用於測試轉換功能的 Word 文件。"
        "包含中英文內容、標題層級與項目符號。"
    )
    doc.add_paragraph(
        "This document is programmatically generated for testing purposes."
    )

    doc.add_heading("第二節 — 功能清單", level=2)
    items = ["支援 PDF 轉換", "支援 HTML 輸出", "支援 Markdown 格式", "多語言支援"]
    for item in items:
        para = doc.add_paragraph(style="List Bullet")
        para.add_run(item)

    doc.add_heading("第三節 — 數據摘要", level=2)
    table = doc.add_table(rows=3, cols=3)
    headers = ["格式", "保真度", "引擎"]
    data = [
        ["PDF", "高", "docx2pdf"],
        ["HTML", "高", "python-docx"],
    ]
    for col_idx, header in enumerate(headers):
        table.cell(0, col_idx).text = header
    for row_idx, row in enumerate(data, start=1):
        for col_idx, cell in enumerate(row):
            table.cell(row_idx, col_idx).text = cell

    doc.save(str(output_path))
    return output_path


def generate_sample_pptx() -> Path:
    """用 python-pptx 建立 3 投影片的 PPTX。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    output_path = FIXTURES_DIR / "sample.pptx"
    prs = Presentation()

    # 投影片 1：標題頁
    slide1_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide1_layout)
    slide1.shapes.title.text = "文件轉檔工具"
    slide1.placeholders[1].text = "測試投影片 — 程式化生成"

    # 投影片 2：內容頁
    slide2_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide2_layout)
    slide2.shapes.title.text = "支援格式"
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "支援的格式清單"
    formats = ["PDF、DOCX、PPTX", "HTML、Markdown", "TXT、圖片（PNG/JPG）"]
    for fmt in formats:
        para = tf.add_paragraph()
        para.text = fmt
        para.level = 1

    # 投影片 3：結尾頁
    slide3_layout = prs.slide_layouts[5]  # 空白
    slide3 = prs.slides.add_slide(slide3_layout)
    from pptx.util import Emu
    txBox = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf3 = txBox.text_frame
    tf3.text = "謝謝使用 / Thank You"

    prs.save(str(output_path))
    return output_path


def generate_sample_html() -> Path:
    """手寫一個包含多種元素的 HTML 文件。"""
    output_path = FIXTURES_DIR / "sample.html"
    html_content = """<!DOCTYPE html>
<html lang="zh-TW">
<head>
  <meta charset="utf-8">
  <title>測試 HTML 文件</title>
  <style>
    body { font-family: "Microsoft JhengHei", Arial, sans-serif; margin: 2em; }
    table { border-collapse: collapse; }
    th, td { border: 1px solid #ccc; padding: 6px 10px; }
    th { background: #f0f0f0; }
  </style>
</head>
<body>
  <h1>測試 HTML 文件</h1>
  <h2>第一節 — 段落</h2>
  <p>這是一段中英文混合的測試文字。This is a test paragraph.</p>
  <p>HTML 文件可以包含 <strong>粗體</strong>、<em>斜體</em> 和 <a href="https://example.com">超連結</a>。</p>

  <h2>第二節 — 清單</h2>
  <ul>
    <li>項目一：PDF 轉換</li>
    <li>項目二：Word 文件</li>
    <li>項目三：Markdown 格式</li>
  </ul>

  <h2>第三節 — 有序清單</h2>
  <ol>
    <li>第一步：選擇來源格式</li>
    <li>第二步：選擇目標格式</li>
    <li>第三步：執行轉換</li>
  </ol>

  <h2>第四節 — 表格</h2>
  <table>
    <tr><th>格式</th><th>說明</th><th>保真度</th></tr>
    <tr><td>PDF</td><td>便攜式文件</td><td>高</td></tr>
    <tr><td>DOCX</td><td>Word 文件</td><td>高</td></tr>
    <tr><td>HTML</td><td>網頁格式</td><td>高</td></tr>
  </table>

  <h2>第五節 — 代碼</h2>
  <pre><code>def hello():
    print("Hello, World!")
    print("你好世界")
</code></pre>
</body>
</html>
"""
    output_path.write_text(html_content, encoding="utf-8")
    return output_path


def generate_sample_md() -> Path:
    """手寫一個包含多種 Markdown 元素的文件。"""
    output_path = FIXTURES_DIR / "sample.md"
    md_content = """# 測試 Markdown 文件

## 第一節 — 段落

這是一段普通段落，包含中英文混合文字。
This is a mixed Chinese-English paragraph for testing.

## 第二節 — 強調語法

**粗體文字** 和 *斜體文字* 的測試。
***粗斜體*** 也應該被正確處理。

## 第三節 — 清單

### 無序清單

- 第一個項目
- 第二個項目
  - 子項目 A
  - 子項目 B
- 第三個項目

### 有序清單

1. 步驟一
2. 步驟二
3. 步驟三

## 第四節 — 連結與代碼

這是一個 [超連結](https://example.com)。

行內代碼：`print("Hello")`

代碼塊：

```python
def convert_file(src, dst):
    print(f"Converting {src} to {dst}")
    return True
```

## 第五節 — 表格

| 格式 | 來源 | 目標 | 保真度 |
|------|------|------|--------|
| PDF  | ✓    | ✓    | 高     |
| DOCX | ✓    | ✓    | 高     |
| HTML | ✓    | ✓    | 高     |

## 第六節 — 引用

> 這是一段引用文字。
> 可以跨多行。

---

文件結束。
"""
    output_path.write_text(md_content, encoding="utf-8")
    return output_path


def generate_sample_txt() -> Path:
    """生成一段中英混合純文字。"""
    output_path = FIXTURES_DIR / "sample.txt"
    txt_content = """文件轉檔工具 — 測試純文字檔案
Document Conversion Tool — Sample Plain Text File

第一段：中文說明
本工具支援多種文件格式之間的轉換，包括 PDF、Word、PowerPoint、HTML、Markdown 等格式。
所有轉換在本機執行，無需連接網路，保護您的隱私安全。

Second Paragraph: English Description
This tool supports conversions between multiple document formats including PDF, Word,
PowerPoint, HTML, Markdown, and plain text. All conversions run locally on your machine
without requiring network access, protecting your privacy.

第三段：功能特點
1. 支援 7 種格式互轉
2. 雙引擎架構（首選 + 降級）
3. 批次轉換支援
4. 美觀的使用者介面

Third Paragraph: Key Features
- No watermarks or usage limits
- Completely offline operation
- Modern drag-and-drop interface
- Automatic format detection

結尾 / End of File
"""
    output_path.write_text(txt_content, encoding="utf-8")
    return output_path


def generate_sample_png() -> Path:
    """用 Pillow 繪製 500×500 測試圖（含文字）。"""
    from PIL import Image, ImageDraw, ImageFont

    output_path = FIXTURES_DIR / "sample.png"
    img = Image.new("RGB", (500, 500), color=(240, 248, 255))  # 淡藍色背景
    draw = ImageDraw.Draw(img)

    # 背景矩形裝飾
    draw.rectangle([10, 10, 490, 490], outline=(100, 150, 200), width=3)

    # 嘗試使用系統字型（Windows），失敗則用預設
    title_font = None
    body_font = None
    try:
        title_font = ImageFont.truetype("C:/Windows/Fonts/msjh.ttc", 28)
        body_font = ImageFont.truetype("C:/Windows/Fonts/msjh.ttc", 18)
    except (IOError, OSError):
        try:
            title_font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 28)
            body_font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 18)
        except (IOError, OSError):
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()

    # 繪製文字
    draw.text((50, 80), "測試圖片", fill=(30, 30, 150), font=title_font)
    draw.text((50, 130), "Test Image", fill=(30, 30, 150), font=title_font)
    draw.text((50, 200), "500 × 500 pixels", fill=(80, 80, 80), font=body_font)
    draw.text((50, 240), "程式化生成 / Programmatically Generated", fill=(80, 80, 80), font=body_font)
    draw.text((50, 290), "用於轉換器測試 / For Converter Testing", fill=(80, 80, 80), font=body_font)

    # 色塊示意
    colors = [
        ((50, 360), (130, 420), (255, 100, 100), "Red"),
        ((160, 360), (240, 420), (100, 255, 100), "Green"),
        ((270, 360), (350, 420), (100, 100, 255), "Blue"),
        ((380, 360), (460, 420), (255, 255, 100), "Yellow"),
    ]
    for (x1, y1), (x2, y2), color, label in colors:
        draw.rectangle([x1, y1, x2, y2], fill=color, outline=(50, 50, 50))

    img.save(str(output_path), "PNG")
    return output_path


def generate_all(force: bool = False) -> dict[str, Path]:
    """生成所有測試 fixtures。

    Args:
        force: 若 True，即使檔案已存在也重新生成

    Returns:
        dict: { fixture_name: path }
    """
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    results: dict[str, Path] = {}

    generators = [
        ("sample.pdf", generate_sample_pdf),
        ("sample.docx", generate_sample_docx),
        ("sample.pptx", generate_sample_pptx),
        ("sample.html", generate_sample_html),
        ("sample.md", generate_sample_md),
        ("sample.txt", generate_sample_txt),
        ("sample.png", generate_sample_png),
    ]

    for filename, gen_func in generators:
        file_path = FIXTURES_DIR / filename
        if force or not file_path.exists():
            try:
                path = gen_func()
                results[filename] = path
                print(f"  生成: {filename}")
            except Exception as exc:
                print(f"  警告: 無法生成 {filename}: {exc}")
        else:
            results[filename] = file_path

    return results


if __name__ == "__main__":
    print("生成測試 fixtures...")
    paths = generate_all(force=True)
    print(f"完成：{len(paths)} 個 fixtures 已生成到 {FIXTURES_DIR}")
